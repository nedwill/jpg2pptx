import os
import sys
from pptx import Presentation
import Image

#take images and turn them into a powerpoint

#checkFile
#if os.path.isdir(savePath):
#		print "We can't save there because a directory exists there."
#		sys.exit(1)
#	if os.path.isfile(savePath):
#		print "File exists."

#checkJPGList

def getDimensions(picWidth,picHeight,slideWidth,slideHeight,emuPixelRatio):
	if (picWidth > slideWidth or picHeight > slideHeight):
		xAdj = picWidth * slideHeight #check these... use multiplication for ratio to avoid introducing rounding issues... check for overflow?
		yAdj = picHeight * slideWidth
		#both less: center
		if xAdj > yAdj:
			#wider than tall
			x = 0 #left side starts at 0
			cx = slideWidth
			scale = float(slideWidth) / picWidth
			centerHeight = slideHeight / 2
			picCenterHeight = picHeight / 2
			y = (centerHeight - (picCenterHeight*scale))
			cy = picHeight*scale
		elif xAdj < yAdj:
			#taller than it is wide
			y = 0
			cy = slideHeight
			scale = float(slideHeight) / picHeight
			centerWidth = slideWidth / 2
			picCenterWidth = picWidth / 2
			x = (centerWidth - (picCenterWidth*scale))
			cx = picWidth*scale
		else:
			#perfectly matches slide ratio, left = 0, top = 0, scale
			assert xAdj == yAdj
			x = 0
			y = 0
			cx = slideWidth
			cy = slideHeight
	else:
		centerWidth = slideWidth / 2
		centerHeight = slideHeight / 2
		picCenterWidth = picWidth / 2
		picCenterHeight = picHeight / 2
		x = centerWidth - picCenterWidth
		cx = picWidth
		y = centerHeight - picCenterHeight
		cy = picHeight
	print (x,y,cx,cy)
	return (int(x*emuPixelRatio),int(y*emuPixelRatio),int(cx*emuPixelRatio),int(cy*emuPixelRatio)) #do this more pythonically?

def jpgListToPowerpoint(jpgList,savePath):
	#slide dimensions are in inches
	slideWidth = 720 #dimensions are @96DPI, default for pptx library
	slideHeight = 540
	emuPixelRatio = 12700
	if not type(jpgList) is list:
		print "jpgList in jpgListToPowerpoint is not a list"
		sys.exit(1)
	prs = Presentation()
	for jpgPath in jpgList:
		blank_slidelayout = prs.slidelayouts[6]
		slide = prs.slides.add_slide(blank_slidelayout)
		(picX,picY) = Image.open(jpgPath).size
		(x, y, cx, cy) = getDimensions(picX,picY,slideWidth,slideHeight,emuPixelRatio)
		#print (picX,picY),(x, y, cx, cy)
		#bug: scaled pictures don't show up in presentation - check documentation
		pic = slide.shapes.add_picture(jpgPath, x, y, cx, cy)
	print "Saving..."
	prs.save(savePath)

#getJPGList

#issues:
	#memory
	#race conditions

if __name__ == '__main__':
	#run cli interface
	pass

#while designing code...
jpgListToPowerpoint(["sample/1.jpg","sample/2.jpg"],"sample/output.pptx")